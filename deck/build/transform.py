#!/usr/bin/env python3
"""Transform Efes Nexus decks (RU + EN):
1. Slide 2: remove auto/human distinction, NEXUS -> SALES SYSTEMS COMMUNICATION SPECIALIST,
   warehouse-availability step -> red (not performed).
2. Slide 4 -> rebuilt as system screenshots slide (was 'Three business processes').
3. Slide 5: merged stats (adds 158h/28% card from old slide 6, no 7000h).
4. Slide 6 -> rebuilt as Future development possibilities (was outlook w/ 7000h).
5. Slide 7 deleted (old 'Next step - scaling').
"""
import copy, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
def q(t):
    pre, loc = t.split(':'); return '{%s}%s' % (NS[pre], loc)
EMU = 914400
def inch(v): return int(round(v * EMU))

# ---------- low-level helpers on lxml sp elements ----------
def sp_texts(sp):
    return ''.join(t.text or '' for t in sp.iter(q('a:t')))

def sp_geom(sp):
    xfrm = sp.find('.//' + q('a:xfrm'))
    if xfrm is None: return None
    off = xfrm.find(q('a:off')); ext = xfrm.find(q('a:ext'))
    if off is None or ext is None: return None
    return (int(off.get('x')), int(off.get('y')), int(ext.get('cx')), int(ext.get('cy')))

def sp_prst(sp):
    g = sp.find('.//' + q('a:prstGeom'))
    return g.get('prst') if g is not None else None

def set_xy(sp, x=None, y=None, w=None, h=None):
    xfrm = sp.find('.//' + q('a:xfrm'))
    off = xfrm.find(q('a:off')); ext = xfrm.find(q('a:ext'))
    if x is not None: off.set('x', str(inch(x)))
    if y is not None: off.set('y', str(inch(y)))
    if w is not None: ext.set('cx', str(inch(w)))
    if h is not None: ext.set('cy', str(inch(h)))

def move_dy(sp, dy):
    off = sp.find('.//' + q('a:off'))
    off.set('y', str(int(off.get('y')) + inch(dy)))

def set_shape_fill(sp, color):
    spPr = sp.find(q('p:spPr'))
    for sf in spPr.findall(q('a:solidFill')):
        c = sf.find(q('a:srgbClr'))
        if c is not None: c.set('val', color)
        return
def set_shape_line(sp, color):
    spPr = sp.find(q('p:spPr'))
    ln = spPr.find(q('a:ln'))
    if ln is None: return
    sf = ln.find(q('a:solidFill'))
    if sf is not None:
        c = sf.find(q('a:srgbClr'))
        if c is not None: c.set('val', color)

def set_run_colors(sp, color):
    for rPr in sp.iter(q('a:rPr')):
        sf = rPr.find(q('a:solidFill'))
        if sf is not None:
            c = sf.find(q('a:srgbClr'))
            if c is not None: c.set('val', color)

def set_run_sizes(sp, sz_pt):
    for rPr in sp.iter(q('a:rPr')):
        rPr.set('sz', str(int(sz_pt * 100)))
        rPr.attrib.pop('spc', None)

def set_single_text(sp, text):
    tx = sp.find(q('p:txBody'))
    paras = tx.findall(q('a:p'))
    p0 = paras[0]
    runs = p0.findall(q('a:r'))
    runs[0].find(q('a:t')).text = text
    for r in runs[1:]: p0.remove(r)
    for pp in paras[1:]: tx.remove(pp)

def set_anchor_top(sp):
    bodyPr = sp.find('.//' + q('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 't')

def delete_sp(sp):
    sp.getparent().remove(sp)

def spTree(slide):
    return slide.shapes._spTree

def all_graphics(slide):
    """direct sp + pic children"""
    tree = spTree(slide)
    return [el for el in tree if el.tag in (q('p:sp'), q('p:pic'))]

# ---------- shadow xml for light cards ----------
SHADOW_XML = ('<a:effectLst xmlns:a="%s"><a:outerShdw sx="100000" sy="100000" kx="0" ky="0" '
              'algn="bl" rotWithShape="0" blurRad="88900" dist="25400" dir="8100000">'
              '<a:srgbClr val="0C1726"><a:alpha val="8000"/></a:srgbClr></a:outerShdw></a:effectLst>') % NS['a']

def add_shadow(shape):
    spPr = shape._element.find(q('p:spPr'))
    spPr.append(etree.fromstring(SHADOW_XML))

def no_shadow(shape):
    spPr = shape._element.find(q('p:spPr'))
    spPr.append(etree.fromstring('<a:effectLst xmlns:a="%s"/>' % NS['a']))

# ---------- python-pptx builders ----------
FONT = 'Calibri'; MONO = 'Consolas'
def C(hexs): return RGBColor.from_string(hexs)

def add_round_rect(slide, x, y, w, h, fill, line=None, radius_in=0.09, shadow=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try: sh.adjustments[0] = radius_in / min(w, h)
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if line: sh.line.color.rgb = C(line); sh.line.width = Pt(1)
    else: sh.line.fill.background()
    if shadow: add_shadow(sh)
    else: no_shadow(sh)
    sh.text_frame.paragraphs[0].text = ''
    return sh

def add_rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    sh.line.fill.background(); no_shadow(sh)
    return sh

def add_text(slide, x, y, w, h, runs, size=12, color='0C1726', bold=False, italic=False,
             align='l', anchor='t', font=FONT, spacing=None, char_spacing=None, wrap=True):
    """runs: str or list of (text, overrides-dict)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}[align]
    if spacing: p.line_spacing = spacing
    if isinstance(runs, str): runs = [(runs, {})]
    for txt, ov in runs:
        r = p.add_run(); r.text = txt
        f = r.font
        f.name = ov.get('font', font); f.size = Pt(ov.get('size', size))
        f.bold = ov.get('bold', bold); f.italic = ov.get('italic', italic)
        f.color.rgb = C(ov.get('color', color))
        if char_spacing or ov.get('char_spacing'):
            r._r.get_or_add_rPr().set('spc', str(int((ov.get('char_spacing', char_spacing)) * 100)))
    return tb

def add_pic(slide, path, x, y, w, h=None):
    return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h) if h else None)

def kicker(slide, x, y, text, dark):
    add_rect(slide, x, y + 0.085, 0.17, 0.032, '5CA3DB' if dark else '004C8D')
    add_text(slide, x + 0.25, y - 0.05, 9, 0.3, text.upper(), size=10.5, bold=True,
             color='9CC6E8' if dark else '1166A8', anchor='m', char_spacing=2)

MX = 0.62; CW = 13.333 - MX * 2

# ============================================================
# SLIDE 2 TRANSFORM
# ============================================================
def transform_slide2(slide, L):
    role_nexus = 'NEXUS'
    delete_texts = {'авто', 'человек', 'auto', 'human',
                    'Автоматизировано системой', 'Действие / решение человека',
                    'Automated by the system', 'Human action / decision'}
    els = all_graphics(slide)
    nodes = []   # (el, x,y,w,h)
    for el in els:
        g = sp_geom(el)
        if g is None: continue
        x, y, w, h = g
        if el.tag == q('p:sp') and sp_prst(el) == 'roundRect' and 1.55*EMU < w < 1.85*EMU and 0.9*EMU < h < 1.1*EMU:
            nodes.append((el, x, y, w, h))
    assert len(nodes) == 13, f'expected 13 nodes, got {len(nodes)}'

    def owner_node(cx, cy):
        for i, (el, x, y, w, h) in enumerate(nodes):
            if x <= cx <= x + w and y - int(0.05*EMU) <= cy <= y + h:
                return i
        return None

    red_idx = None
    node_children = {i: [] for i in range(len(nodes))}
    to_delete = []
    for el in els:
        g = sp_geom(el)
        if g is None: continue
        x, y, w, h = g
        cx, cy = x + w // 2, y + h // 2
        txt = sp_texts(el) if el.tag == q('p:sp') else ''
        top_in = y / EMU

        # legend + callout zone
        if 6.25 < top_in < 6.9:
            to_delete.append(el); continue
        # subtitle
        if txt.startswith('Шаги BPMN') or txt.startswith('Steps of the BPMN'):
            newsub = ('Шаги BPMN-диаграммы и ответственные роли. Красным выделен шаг, который сегодня никем не выполняется.'
                      if L == 'ru' else
                      'Steps of the BPMN diagram and the responsible roles. Red marks the step that no one performs today.')
            set_single_text(el, newsub)
            continue
        # arrows between nodes
        if txt == '→' and 0.9*EMU < h < 1.1*EMU:
            set_run_colors(el, 'B9C1CC'); continue

        ni = owner_node(cx, cy)
        if ni is None: continue
        if el in [n[0] for n in nodes]: continue
        # inside a node
        if el.tag == q('p:pic'):
            to_delete.append(el); continue
        if txt in delete_texts:
            to_delete.append(el); continue
        prst = sp_prst(el)
        if prst == 'ellipse' and w < 0.12*EMU:   # human dot
            to_delete.append(el); continue
        node_children[ni].append((el, x, y, w, h, txt, prst))
        if ('склад' in txt) or ('warehouse' in txt):
            red_idx = ni

    for el in to_delete: delete_sp(el)
    assert red_idx is not None, 'red node not found'

    ROLES = {'NEXUS', 'TRADE-МАРКЕТИНГ', 'CRM-СПЕЦИАЛИСТ', 'TM / ON-TRADE', 'ПАРТНЁР',
             'TRADE MARKETING', 'CRM SPECIALIST', 'PARTNER'}
    for i, (nel, nx, ny, nw, nh) in enumerate(nodes):
        red = (i == red_idx)
        set_shape_fill(nel, 'FBE9E7' if red else 'FFFFFF')
        set_shape_line(nel, 'F0C5C0' if red else 'E2E7EE')
        for (el, x, y, w, h, txt, prst) in node_children[i]:
            if prst == 'rect' and w < 0.09*EMU:          # accent bar
                set_shape_fill(el, 'C33A30' if red else '004C8D')
            elif prst == 'ellipse':                       # number circle
                set_shape_fill(el, 'C33A30' if red else '004C8D')
            elif txt and len(txt) == 1 and txt.isdigit():  # number
                set_run_colors(el, 'FFFFFF')
            elif txt in ROLES:                            # role label
                if red:
                    newrole = 'НЕ ВЫПОЛНЯЕТСЯ' if L == 'ru' else 'NOT PERFORMED'
                    set_single_text(el, newrole); set_run_colors(el, 'C33A30')
                elif txt == role_nexus:
                    set_single_text(el, 'SALES SYSTEMS COMMUNICATION SPECIALIST')
                    set_run_colors(el, '1166A8')
                else:
                    set_run_colors(el, '1166A8')
                set_run_sizes(el, 7)
                set_xy(el, x=nx/EMU + 0.46, y=ny/EMU + 0.10, w=nw/EMU - 0.54, h=0.42)
                set_anchor_top(el)
            else:                                         # step name
                set_xy(el, y=ny/EMU + 0.545, h=0.42)
                set_anchor_top(el)

# ============================================================
# STRIP SLIDE (keep footer & background)
# ============================================================
def strip_slide(slide):
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g is None: continue
        if g[1] < 6.9 * EMU:
            delete_sp(el)

# ============================================================
# SLIDE 4 -> SCREENSHOTS
# ============================================================
def build_screenshots(slide, L, shots):
    strip_slide(slide)
    S = {
      'ru': dict(kick='Система', h2='Промо-дашборд',
                 quote='«Единое окно для всего жизненного цикла промо-активности — от заявки партнёра до расчёта компенсации.»',
                 cap1b='Реестр промо-активностей', cap1=' — фильтры, квоты и реализация в реальном времени',
                 cap2b='Карточка промо-активности', cap2=' — охват, точки продаж, продукты и операции'),
      'en': dict(kick='The system', h2='The promo dashboard',
                 quote='“A single window for the entire promo-activity lifecycle — from the partner’s request to the compensation calculation.”',
                 cap1b='Promo-activity registry', cap1=' — filters, quotas and realisation in real time',
                 cap2b='Promo-activity details', cap2=' — coverage, sale points, products and operations'),
    }[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_pic(slide, str(ICONDIR / 'quote.png'), MX + 0.02, 1.70, 0.24, 0.24)
    add_text(slide, MX + 0.4, 1.62, CW - 0.5, 0.6, S['quote'], size=14.5, italic=True,
             color='13203A', spacing=1.2)
    cw = 5.85; gap = CW - 2 * cw
    for k, (img, capb, cap) in enumerate([(shots[0], S['cap1b'], S['cap1']), (shots[1], S['cap2b'], S['cap2'])]):
        x = MX + k * (cw + gap)
        add_round_rect(slide, x, 2.42, cw, 3.99, 'FFFFFF', 'E2E7EE', shadow=True)
        pic = add_pic(slide, img, x + 0.12, 2.54, cw - 0.24, (cw - 0.24) / (3200/1900))
        pic.line.color.rgb = C('E2E7EE'); pic.line.width = Pt(0.75)
        add_text(slide, x + 0.16, 5.95, cw - 0.32, 0.4,
                 [(capb, {'bold': True, 'color': '0C1726'}), (cap, {'color': '56616F'})],
                 size=11, anchor='m')

# ============================================================
# SLIDE 5 -> MERGED STATS
# ============================================================
def merge_stats(slide, L):
    els = all_graphics(slide)
    right = [el for el in els if (g := sp_geom(el)) and g[0] > 6.5*EMU and 2.5*EMU < g[1] < 6.4*EMU]
    rows = {1: [], 2: [], 3: []}
    for el in right:
        top = sp_geom(el)[1] / EMU
        if top < 3.7: rows[1].append(el)
        elif top < 5.0: rows[2].append(el)
        else: rows[3].append(el)
    # clone row3 as template for new card BEFORE moving
    tmpl = {}
    for el in rows[3]:
        g = sp_geom(el); txt = sp_texts(el); prst = sp_prst(el)
        key = None
        if prst == 'roundRect' and g[2] > 5*EMU: key = 'card'
        elif prst == 'roundRect' and g[2] < 1.2*EMU and not txt: key = 'pillrect'
        elif txt.startswith('СУММА') or txt.startswith('PROMO TOTAL'): key = 'label'
        elif '2023' in txt: key = 'rng'
        elif txt in ('2,41', '2.41'): key = 'b'
        elif '+140' in txt: key = 'pilltext'
        if key: tmpl[key] = el
    # move rows: targets 2.54 / 3.56 / 4.58 (deltas -0.25, -0.51, -0.77)
    for rowi, dy in ((1, -0.25), (2, -0.51), (3, -0.77)):
        for el in rows[rowi]: move_dy(el, dy)
    # build 4th card at 5.60 by cloning template (orig row3 base was 5.35 pre-move)
    T = {'ru': dict(label='ВОЗВРАЩЁННОЕ РАБОЧЕЕ ВРЕМЯ', rng='Год 1 · автоматизация',
                    val='158 ч/мес', pill='28 % рутины — у системы'),
         'en': dict(label='WORKING TIME RETURNED', rng='Year 1 · automation',
                    val='158 h/mo', pill='28% of routine automated')}[L]
    tree = spTree(slide)
    dy = 1.02  # template row already moved to 4.58; new card sits at 5.60
    for key in ('card', 'label', 'rng', 'b', 'pillrect', 'pilltext'):
        el = copy.deepcopy(tmpl[key]); tree.append(el); move_dy(el, dy)
        if key == 'label': set_single_text(el, T['label'])
        if key == 'rng':
            set_single_text(el, T['rng'])
            g = sp_geom(el); set_xy(el, x=g[0]/EMU - 0.6, w=g[2]/EMU + 0.6)
        if key == 'b':
            set_single_text(el, T['val'])
            set_xy(el, x=7.24, w=2.6)
        if key in ('pillrect', 'pilltext'):
            g = sp_geom(el)
            neww = 2.05
            set_xy(el, x=g[0]/EMU - (neww - g[2]/EMU), w=neww)
            if key == 'pilltext': set_single_text(el, T['pill'])

# ============================================================
# SLIDE 6 -> FUTURE DEVELOPMENT
# ============================================================
def build_future(slide, L):
    strip_slide(slide)
    S = {
     'ru': dict(kick='Развитие', h2='Этап 2 · Развитие платформы',
        sub='Платформа уже в работе — опциональный Этап 2 (5 000 000 ₸) расширяет её: AI-поиск по истории акций, проверка накладных 1С, новые типы расчётов и календарь промо.',
        cards=[
          dict(t='AI Search & Knowledge Base', p='1 600 000 ₸',
               b='Поиск и ответы на вопросы по текущим и историческим промо-акциям.',
               lead='Пример:', ex='«Какие акции были на Ульбу в июле 2024, 2025, 2026?»'),
          dict(t='1C Invoice Validation', p='1 500 000 ₸',
               b='Автоматическая проверка выгрузок и накладных из 1С.',
               lead='Эффект:', ex='−30 % времени трейд- и sales controlling-специалистов; проверка 1С каждой точки вручную — ∞ часов.'),
          dict(t='Исторический календарь промо', p='800 000 ₸',
               b='Календарь всех текущих и исторических промо-акций.',
               lead='Как работает:', ex='клик по дате показывает активности, действовавшие в этот период.'),
          dict(t='Масштабирование расчётов', p='1 100 000 ₸',
               b='Перенос расчётов на другие типы компенсаций — например, расчёт Kega Boom.',
               lead='База готова:', ex='единая логика квот, цен и периодов уже в системе.'),
        ]),
     'en': dict(kick="What's next", h2='Phase 2 · Platform development',
        sub='The platform is already live — the optional Phase 2 (5,000,000 KZT) extends it: AI search over promo history, 1C invoice validation, new calculation types and a promo calendar.',
        cards=[
          dict(t='AI Search & Knowledge Base', p='1,600,000 KZT',
               b='Search and answers to questions about current and historical promo activities.',
               lead='Example:', ex='“Which promos ran for Ulba in July 2024, 2025 and 2026?”'),
          dict(t='1C Invoice Validation', p='1,500,000 KZT',
               b='Automatic validation of exports and invoices from 1C.',
               lead='Effect:', ex='−30% of trade & sales-controlling specialists’ time; checking 1C from every point manually — ∞ hours.'),
          dict(t='Historical Promo Calendar', p='800,000 KZT',
               b='A calendar of all current and historical promo activities.',
               lead='How it works:', ex='clicking a date shows the activities that were active in that period.'),
          dict(t='Scaling the calculations', p='1,100,000 KZT',
               b='Extending calculations to other compensation types — e.g. the Kega Boom calculation.',
               lead='Foundation ready:', ex='shared quota, price and period logic is already in the system.'),
        ]),
    }[L]
    ICONS = [(str(ICONDIR / 'search.png'), '14304F'), (str(ICONDIR / 'invoice.png'), '2C2310'),
             (str(ICONDIR / 'calendar.png'), '14304F'), (str(ICONDIR / 'calc.png'), '10301F')]
    kicker(slide, MX, 0.56, S['kick'], True)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='EEF2F7', anchor='m')
    add_text(slide, MX, 1.56, 11.9, 0.8, S['sub'], size=12.5, color='9AA9BF', spacing=1.15)
    gx = 0.28; cw = (CW - gx * 3) / 4; cy = 2.42; ch = 3.62
    for i, cardS in enumerate(S['cards']):
        x = MX + i * (cw + gx)
        add_round_rect(slide, x, cy, cw, ch, '16223B', '28364E')
        chip, icocolor = ICONS[i]
        d = 0.58
        add_round_rect(slide, x + 0.28, cy + 0.3, d, d, ICONS[i][1], None, radius_in=d*0.27)
        add_pic(slide, chip, x + 0.28 + d*0.24, cy + 0.3 + d*0.24, d*0.52, d*0.52)
        # module price chip (merged from the other machine's Этап 2 slide)
        pw2 = 1.34
        add_round_rect(slide, x + cw - 0.28 - pw2, cy + 0.42, pw2, 0.34, '0E1726', '2E72AE',
                       radius_in=0.17)
        add_text(slide, x + cw - 0.28 - pw2, cy + 0.42, pw2, 0.34, cardS['p'], size=9.5,
                 bold=True, color='9CC6E8', align='c', anchor='m')
        add_text(slide, x + 0.28, cy + 1.02, cw - 0.56, 0.72, cardS['t'], size=14, bold=True,
                 color='EEF2F7', spacing=1.05)
        add_text(slide, x + 0.28, cy + 1.78, cw - 0.56, 0.85, cardS['b'], size=10.5, color='9AA9BF', spacing=1.16)
        add_round_rect(slide, x + 0.22, cy + 2.62, cw - 0.44, 0.88, '0E1726', '28364E', radius_in=0.07)
        add_text(slide, x + 0.36, cy + 2.72, cw - 0.72, 0.72,
                 [(cardS['lead'] + ' ', {'bold': True, 'color': '5CA3DB', 'italic': True}),
                  (cardS['ex'], {'italic': True, 'color': '9AA9BF'})],
                 size=9.5, spacing=1.1)
    # accent band: new promo mechanics + instant partner notifications
    by = 6.22; bh = 0.58
    add_round_rect(slide, MX, by, CW, bh, '0C3F77', '2E72AE', radius_in=0.09)
    add_pic(slide, str(ICONDIR / 'bolt_white.png'), MX + 0.32, by + bh/2 - 0.13, 0.26, 0.26)
    band = ('Новые механики легко интегрируются в платформу. Полный пакет развития (Этап 2) — '
            '5 000 000 ₸; модули можно запускать по отдельности.'
            if L == 'ru' else
            'New mechanics plug straight into the platform. The full development package (Phase 2) — '
            '5,000,000 KZT; modules can be started separately.')
    add_text(slide, MX + 0.76, by, CW - 1.1, bh, band, size=12.5, bold=True, color='FFFFFF', anchor='m')

# ============================================================
# helpers for node-group work on BPMN slides
# ============================================================
FNW = (CW - 0.34 * 5) / 6  # node width
ARW = 0.34

def move_dx(sp, dx):
    off = sp.find('.//' + q('a:off'))
    off.set('x', str(int(off.get('x')) + inch(dx)))

def find_nodes(slide):
    out = []
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g is None: continue
        x, y, w, h = g
        if el.tag == q('p:sp') and sp_prst(el) == 'roundRect' and 1.55*EMU < w < 1.85*EMU and 0.9*EMU < h < 1.1*EMU:
            out.append((el, x, y, w, h))
    return out

def node_children_map(slide, nodes):
    kids = {i: [] for i in range(len(nodes))}
    for el in all_graphics(slide):
        if el in [n[0] for n in nodes]: continue
        g = sp_geom(el)
        if g is None: continue
        x, y, w, h = g
        cx, cy = x + w // 2, y + h // 2
        for i, (nel, nx, ny, nw, nh) in enumerate(nodes):
            if nx <= cx <= nx + nw and ny - int(0.05*EMU) <= cy <= ny + nh:
                kids[i].append(el); break
    return kids

# ============================================================
# 1) remove the red (not-performed) step from slide 2 row 01
# ============================================================
def remove_red_node(slide, L):
    nodes = find_nodes(slide)
    kids = node_children_map(slide, nodes)
    row1 = sorted([i for i, n in enumerate(nodes) if n[2] < 3.0*EMU],
                  key=lambda i: nodes[i][1])
    assert len(row1) == 6
    # red node = the one whose role text says not-performed
    red_i = None
    for i in row1:
        for el in kids[i]:
            t = sp_texts(el)
            if t in ('НЕ ВЫПОЛНЯЕТСЯ', 'NOT PERFORMED'): red_i = i
    assert red_i is not None
    # delete red node + children
    delete_sp(nodes[red_i][0])
    for el in kids[red_i]: delete_sp(el)
    # arrows in row 1
    arrows = []
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g is None: continue
        if sp_texts(el) == '→' and 0.9*EMU < g[3] < 1.1*EMU and g[1] < 3.4*EMU:
            arrows.append((el, g[0]))
    arrows.sort(key=lambda t: t[1])
    assert len(arrows) == 5
    delete_sp(arrows[0][0])  # arrow before the (deleted) 2nd node
    arrows = arrows[1:]
    # reposition remaining 5 nodes centered
    remaining = [i for i in row1 if i != red_i]
    remaining.sort(key=lambda i: nodes[i][1])
    rowW = 5 * FNW + 4 * ARW
    startX = MX + (CW - rowW) / 2
    for k, i in enumerate(remaining):
        nel, nx, ny, nw, nh = nodes[i]
        dx = (startX + k * (FNW + ARW)) - nx / EMU
        move_dx(nel, dx)
        for el in kids[i]:
            move_dx(el, dx)
            t = sp_texts(el)
            if len(t) == 1 and t.isdigit():
                set_single_text(el, str(k + 1))
    for k, (ael, ax) in enumerate(arrows):
        set_xy(ael, x=startX + k * (FNW + ARW) + FNW, w=ARW)
    # subtitle without the red mention
    for el in all_graphics(slide):
        t = sp_texts(el)
        if t.startswith('Шаги BPMN-диаграммы и') or t.startswith('Steps of the BPMN diagram and'):
            set_single_text(el, 'Шаги BPMN-диаграммы и ответственные роли — как процесс выглядит сегодня.'
                            if L == 'ru' else
                            'Steps of the BPMN diagram and the responsible roles — how the process works today.')

# ============================================================
# 3+10) mark warehouse step as NEW (red) on the with-Nexus slide,
#       replace the callout with the "not just automation" accent
# ============================================================
def mark_new_step(slide, L):
    tree = slide.part._element.find(q('p:cSld')).find(q('p:spTree'))
    els = [el for el in tree if el.tag in (q('p:sp'), q('p:pic'))]
    def geoms(el): return sp_geom(el)
    nodes = []
    for el in els:
        g = sp_geom(el)
        if g and el.tag == q('p:sp') and sp_prst(el) == 'roundRect' and 1.55*EMU < g[2] < 1.85*EMU and 0.9*EMU < g[3] < 1.1*EMU:
            nodes.append((el, *g))
    # children map (local, since slide.shapes may be stale)
    kids = {i: [] for i in range(len(nodes))}
    for el in els:
        if el in [n[0] for n in nodes]: continue
        g = sp_geom(el)
        if g is None: continue
        cx, cy = g[0] + g[2]//2, g[1] + g[3]//2
        for i, (nel, nx, ny, nw, nh) in enumerate(nodes):
            if nx <= cx <= nx + nw and ny - int(0.05*EMU) <= cy <= ny + nh:
                kids[i].append(el); break
    tgt = None
    for i in range(len(nodes)):
        if any(('склад' in sp_texts(el)) or ('warehouse' in sp_texts(el)) for el in kids[i]):
            tgt = i
    assert tgt is not None
    nel = nodes[tgt][0]
    set_shape_fill(nel, 'FBE9E7'); set_shape_line(nel, 'F0C5C0')
    for el in list(kids[tgt]):
        t = sp_texts(el); prst = sp_prst(el); g = sp_geom(el)
        if el.tag == q('p:pic'):
            delete_sp(el)  # blue check icon
        elif prst == 'rect' and g[2] < 0.09*EMU:
            set_shape_fill(el, 'C33A30')
        elif prst == 'ellipse':
            set_shape_fill(el, 'C33A30')
        elif t == 'NEXUS':
            set_run_colors(el, 'C33A30')
        elif t in ('авто', 'auto'):
            set_single_text(el, 'новый шаг' if L == 'ru' else 'new step')
            set_run_colors(el, 'C33A30')
            set_xy(el, x=g[0]/EMU - 0.16)
    # subtitle: append red note
    for el in els:
        t = sp_texts(el)
        if t.startswith('Шаги BPMN-диаграммы.') or t.startswith('Steps of the BPMN diagram.'):
            tx = el.find(q('p:txBody')); p0 = tx.findall(q('a:p'))[0]
            runs = p0.findall(q('a:r'))
            for txt, red, bold in ((' ', False, False),
                                   (('Красным' if L=='ru' else 'Red'), True, True),
                                   ((' — новый шаг, которого сегодня нет.' if L=='ru'
                                     else ' — a new step that does not exist today.'), False, False)):
                nr = copy.deepcopy(runs[0])
                nr.find(q('a:t')).text = txt
                rPr = nr.find(q('a:rPr'))
                if red:
                    sf = rPr.find(q('a:solidFill'))
                    if sf is not None: sf.find(q('a:srgbClr')).set('val', 'C33A30')
                    rPr.set('b', '1')
                p0.append(nr)
    # legend: third (red) item — clone first swatch + label
    swatch = label = None
    for el in els:
        g = sp_geom(el)
        if g is None: continue
        if 6.3*EMU < g[1] < 6.75*EMU:
            t = sp_texts(el)
            if sp_prst(el) == 'rect' and g[2] < 0.3*EMU and swatch is None:
                swatch = el
            elif t.startswith('Автоматизировано') or t.startswith('Automated'):
                label = el
    if swatch is not None and label is not None:
        s2 = copy.deepcopy(swatch); tree.append(s2)
        set_xy(s2, x=5.95); set_shape_fill(s2, 'FBE9E7'); set_shape_line(s2, 'F0C5C0')
        l2 = copy.deepcopy(label); tree.append(l2)
        set_xy(l2, x=6.25, w=2.2)
        set_single_text(l2, 'Новый шаг процесса' if L == 'ru' else 'New process step')
    # callout -> "not just automation" accent
    for el in els:
        g = sp_geom(el)
        if g is None: continue
        if 6.25*EMU < g[1] < 6.75*EMU and g[0] > 7.0*EMU:
            if sp_prst(el) == 'roundRect':
                set_xy(el, y=6.30, h=0.60)
            elif '8 ' in sp_texts(el):
                set_xy(el, y=6.30, h=0.60)
                tx = el.find(q('p:txBody')); p0 = tx.findall(q('a:p'))[0]
                runs = p0.findall(q('a:r'))
                r1, r2 = runs[0], runs[1]
                lead = 'Не просто автоматизация:' if L == 'ru' else 'Not just automation:'
                rest = (' добавляются новые шаги — проверка остатков на складе партнёра и предрасчёты промо, чтобы партнёры присылали корректные запросы.'
                        if L == 'ru' else
                        ' new steps are added — partner stock checks and promo pre-calculations, so partners send correct requests.')
                r1.find(q('a:t')).text = lead
                lat = r1.find(q('a:rPr')).find(q('a:latin'))
                if lat is not None: lat.set('typeface', FONT)
                r2.find(q('a:t')).text = rest
                for r in runs[2:]: p0.remove(r)

# ============================================================
# 2c) stats slide v3 — no charts; saved time is THE accent
# ============================================================
def stats_v3(slide, L):
    strip_slide(slide)
    S = {'ru': dict(
            kick='Результат в цифрах', h2='Главный эффект — возвращённое время',
            sub='Объём промо вырос почти в 5 раз при той же команде. Платформа возвращает специалистам рабочее время, забирая рутину на себя.',
            big='158 ч/мес', lab='возвращённого рабочего времени · Год 1',
            pill='28 % рутины — у системы', fte='≈ 1,1 штатной ставки, освобождённой от ручной рутины',
            note='Платформа сама переносит данные между Excel, Panorama и письмами партнёрам. Это время специалисты возвращают в работу с партнёрами и в контроль расчётов.',
            tiles=[('ПРОМО-АКТИВНОСТИ / ГОД', '×4,6', '2 187 → 10 112', '+362 %'),
                   ('ПАРТНЁРЫ', '33', '31 → 33 · за 6 мес', '+6,5 %'),
                   ('ПРОДУКТЫ · SKU', '105', '78 → 105 · за 6 мес', '+35 %'),
                   ('СУММА ПРОМО / ГОД', '2,41 млрд ₸', '1,00 → 2,41 · 2023–2025', '+140 %')]),
         'en': dict(
            kick='The result in numbers', h2='The main effect — time returned',
            sub='Promo volume grew almost 5× with the same team. The platform gives specialists their working time back by taking over the routine.',
            big='158 h/mo', lab='of working time returned · Year 1',
            pill='28% of routine automated', fte='≈ 1.1 FTE freed from manual routine',
            note='The platform moves data between Excel, Panorama and partner emails by itself. Specialists put that time back into working with partners and checking calculations.',
            tiles=[('PROMOS / YEAR', '×4.6', '2,187 → 10,112', '+362%'),
                   ('PARTNERS', '33', '31 → 33 · in 6 mo', '+6.5%'),
                   ('PRODUCTS · SKU', '105', '78 → 105 · in 6 mo', '+35%'),
                   ('PROMO TOTAL / YEAR', '2.41 bln ₸', '1.00 → 2.41 · 2023–2025', '+140%')])}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_text(slide, MX, 1.56, 11.9, 0.6, S['sub'], size=12.5, color='56616F', spacing=1.15)
    top, bot = 2.32, 6.46
    # ---- supporting KPI tiles: one row across the top ----
    th = 1.42; gap = 0.20
    tw = (CW - 3 * gap) / 4
    for i, (lab, big, sub, delta) in enumerate(S['tiles']):
        x = MX + i * (tw + gap)
        add_round_rect(slide, x, top, tw, th, 'FFFFFF', 'E2E7EE', shadow=True)
        add_rect(slide, x, top + 0.02, 0.05, th - 0.04, '004C8D')
        add_text(slide, x + 0.30, top + 0.15, tw - 0.5, 0.24, lab, size=8.5, bold=True,
                 color='56616F', char_spacing=0.8, anchor='m')
        add_text(slide, x + 0.30, top + 0.44, tw - 0.5, 0.46, big, size=23, bold=True,
                 color='004C8D', anchor='m')
        add_text(slide, x + 0.30, top + 0.97, tw - 0.5, 0.32,
                 [(sub + '  ', {'color': '6E7885'}), (delta, {'bold': True, 'color': '15924F'})],
                 size=9.5, anchor='m')
    # ---- hero card: the returned hours, full width, below the tiles ----
    hy = top + th + 0.26; hh = bot - hy
    add_round_rect(slide, MX, hy, CW, hh, 'FFFFFF', 'E2E7EE', shadow=True)
    add_rect(slide, MX + 0.02, hy, CW - 0.04, 0.05, '15924F')
    # clock chip, vertically centred on the big number
    ncy = hy + 0.55 + 0.575                       # centre of the 158 h/mo block
    add_round_rect(slide, MX + 0.62, ncy - 0.41, 0.82, 0.82, 'E7F5EC', 'BBE3C8', radius_in=0.22)
    add_pic(slide, str(ICONDIR / 'clock_dkgreen.png'), MX + 0.82, ncy - 0.21, 0.42, 0.42)
    add_text(slide, MX + 1.74, hy + 0.55, 4.7, 1.15, S['big'], size=62, bold=True,
             color='15924F', anchor='m')
    add_text(slide, MX + 1.78, hy + 1.72, 4.7, 0.40, S['lab'], size=14.5, color='56616F', anchor='m')
    # divider + supporting column on the right
    add_rect(slide, MX + 6.62, hy + 0.48, 0.014, hh - 0.96, 'E2E7EE')
    rx = MX + 7.12; rw = CW - 7.12 - 0.55
    ry = hy + (hh - 1.52) / 2
    ppw = 2.42
    add_round_rect(slide, rx, ry, ppw, 0.40, 'E7F5EC', 'BBE3C8', radius_in=0.2)
    add_text(slide, rx, ry, ppw, 0.40, S['pill'], size=11.5, bold=True,
             color='15924F', align='c', anchor='m')
    add_text(slide, rx, ry + 0.58, rw, 0.94, S['note'], size=12.5, color='56616F', spacing=1.25)

# ============================================================
# 2b) stats slide v2 — full rebuild: KPI strip + chart + hero tile
# ============================================================
def stats_v2(slide, L):
    strip_slide(slide)
    S = {'ru': dict(
            kick='Масштаб', h2='Объём растёт быстрее, чем успевают руки',
            sub='За три года число промо-активностей выросло почти в 5 раз при почти неизменном числе партнёров — ручные процессы такого темпа не выдерживают.',
            tiles=[('ПРОМО-АКТИВНОСТИ / ГОД', '×4,6', '2 187 → 10 112', '+362 %'),
                   ('ПАРТНЁРЫ', '33', '31 → 33 · за 6 мес', '+6,5 %'),
                   ('ПРОДУКТЫ · SKU', '105', '78 → 105 · за 6 мес', '+35 %'),
                   ('СУММА ПРОМО / ГОД', '2,41 млрд ₸', '1,00 → 2,41 · 2023–2025', '+140 %')],
            ct='Промо-активности по годам', cs='Количество заведённых акций за год',
            vals=['2 084', '2 187', '10 112'],
            big='158 ч/мес', lab='возвращённого рабочего времени · Год 1',
            pill='28 % рутины — у системы', fte='≈ 1,1 штатной ставки — без ручной рутины'),
         'en': dict(
            kick='Scale', h2='Volume grows faster than hands can keep up',
            sub='In three years the number of promo activities grew almost 5× with an almost unchanged number of partners — manual processes cannot keep this pace.',
            tiles=[('PROMOS / YEAR', '×4.6', '2,187 → 10,112', '+362%'),
                   ('PARTNERS', '33', '31 → 33 · in 6 mo', '+6.5%'),
                   ('PRODUCTS · SKU', '105', '78 → 105 · in 6 mo', '+35%'),
                   ('PROMO TOTAL / YEAR', '2.41 bln ₸', '1.00 → 2.41 · 2023–2025', '+140%')],
            ct='Promo activities by year', cs='Number of promos created per year',
            vals=['2,084', '2,187', '10,112'],
            big='158 h/mo', lab='of working time returned · Year 1',
            pill='28% of routine automated', fte='≈ 1.1 FTE freed from manual routine')}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_text(slide, MX, 1.56, 11.9, 0.6, S['sub'], size=12.5, color='56616F', spacing=1.15)
    # ---- KPI strip ----
    tw = (CW - 3 * 0.17) / 4; ty = 2.32; th = 1.02
    for i, (lab, big, sub, delta) in enumerate(S['tiles']):
        x = MX + i * (tw + 0.17)
        add_round_rect(slide, x, ty, tw, th, 'FFFFFF', 'E2E7EE', shadow=True)
        add_rect(slide, x + 0.02, ty, tw - 0.04, 0.045, '004C8D')
        add_text(slide, x + 0.22, ty + 0.14, tw - 0.4, 0.22, lab, size=8.5, bold=True,
                 color='56616F', char_spacing=0.8, anchor='m')
        add_text(slide, x + 0.22, ty + 0.36, tw - 0.4, 0.4, big, size=22, bold=True, color='004C8D', anchor='m')
        add_text(slide, x + 0.22, ty + 0.74, tw - 0.4, 0.22,
                 [(sub + '  ', {'color': '8B94A2'}), (delta, {'bold': True, 'color': '15924F'})],
                 size=9, anchor='m')
    # ---- chart card ----
    cx, cy, cw, ch = MX, 3.60, 7.06, 2.86
    add_round_rect(slide, cx, cy, cw, ch, 'FFFFFF', 'E2E7EE', shadow=True)
    add_text(slide, cx + 0.3, cy + 0.2, 4.5, 0.3, S['ct'], size=13.5, bold=True, color='0C1726', anchor='m')
    add_text(slide, cx + 0.3, cy + 0.5, 4.5, 0.24, S['cs'], size=9.5, color='8B94A2', anchor='m')
    baseY = cy + ch - 0.52; maxH = 1.52
    years = ['2023', '2024', '2025']; raw = [2084, 2187, 10112]
    bw = 0.95
    add_rect(slide, cx + 0.3, baseY, cw - 0.6, 0.012, 'E2E7EE')
    for i in range(3):
        bh = maxH * raw[i] / raw[2]
        bx = cx + 0.85 + i * 2.05
        last = i == 2
        add_rect(slide, bx, baseY - bh, bw, bh, '004C8D' if last else 'C5CDD8')
        if last: add_rect(slide, bx, baseY - bh, bw, 0.26, '2D83C4')
        add_text(slide, bx - 0.35, baseY - bh - 0.32, bw + 0.7, 0.28, S['vals'][i],
                 size=12 if not last else 13, bold=True, color='004C8D' if last else '56616F',
                 font=MONO, align='c', anchor='m')
        add_text(slide, bx - 0.35, baseY + 0.08, bw + 0.7, 0.26, years[i], size=11, bold=True,
                 color='004C8D' if last else '8B94A2', align='c', anchor='m')
    # ---- hero tile (same white card language, green accent) ----
    hx = MX + cw + 0.18; hw = CW - cw - 0.18
    add_round_rect(slide, hx, cy, hw, ch, 'FFFFFF', 'E2E7EE', shadow=True)
    add_rect(slide, hx + 0.02, cy, hw - 0.04, 0.045, '15924F')
    add_round_rect(slide, hx + 0.34, cy + 0.34, 0.56, 0.56, 'E7F5EC', 'BBE3C8', radius_in=0.16)
    add_pic(slide, str(ICONDIR / 'clock_dkgreen.png'), hx + 0.34 + 0.13, cy + 0.34 + 0.13, 0.30, 0.30)
    add_text(slide, hx + 1.06, cy + 0.32, hw - 1.3, 0.6, S['big'], size=40, bold=True, color='15924F', anchor='m')
    add_text(slide, hx + 0.36, cy + 1.08, hw - 0.7, 0.5, S['lab'], size=12.5, color='56616F', spacing=1.15)
    add_rect(slide, hx + 0.36, cy + 1.70, hw - 0.72, 0.014, 'E2E7EE')
    ppw = 2.15
    add_round_rect(slide, hx + 0.36, cy + 1.92, ppw, 0.34, 'E7F5EC', 'BBE3C8', radius_in=0.17)
    add_text(slide, hx + 0.36, cy + 1.92, ppw, 0.34, S['pill'], size=10.5, bold=True,
             color='15924F', align='c', anchor='m')
    add_text(slide, hx + 0.36, cy + 2.38, hw - 0.7, 0.3, S['fte'], size=11, color='8B94A2', anchor='m')

# ============================================================
# 2) redesigned stats slide: compact metric rows + big accent tile
# ============================================================
def redesign_stats(slide, L):
    # remove old right column
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g is None: continue
        if g[0] > 6.5*EMU and 2.4*EMU < g[1] < 6.5*EMU:
            delete_sp(el)
    R = {'ru': [('ПАРТНЁРЫ', 'за 6 мес', '31', '33', '+6,5 %'),
                ('ПРОДУКТЫ · SKU', 'за 6 мес', '78', '105', '+35 %'),
                ('СУММА ПРОМО / ГОД · МЛРД ₸', '2023 → 2025', '1,00', '2,41', '+140 %')],
         'en': [('PARTNERS', 'in 6 mo', '31', '33', '+6.5%'),
                ('PRODUCTS · SKU', 'in 6 mo', '78', '105', '+35%'),
                ('PROMO TOTAL / YEAR · BLN ₸', '2023 → 2025', '1.00', '2.41', '+140%')]}[L]
    X, W = 6.96, 5.75
    for i, (lab, rng, a, b, d) in enumerate(R):
        y = 2.54 + i * 0.90
        add_round_rect(slide, X, y, W, 0.76, 'FFFFFF', 'E2E7EE', shadow=True)
        add_text(slide, X + 0.28, y + 0.10, 3.9, 0.26, lab, size=9.5, bold=True, color='56616F', anchor='m', char_spacing=1)
        add_text(slide, X + W - 2.85, y + 0.10, 1.5, 0.26, rng, size=9, color='8B94A2', font=MONO, align='r', anchor='m')
        add_text(slide, X + 0.28, y + 0.34, 1.1, 0.36, a, size=14, bold=True, color='8B94A2', font=MONO, anchor='m')
        add_text(slide, X + 1.35, y + 0.34, 0.5, 0.36, '→', size=13, bold=True, color='5CA3DB', align='c', anchor='m')
        add_text(slide, X + 1.95, y + 0.31, 1.6, 0.4, b, size=19, bold=True, color='004C8D', font=MONO, anchor='m')
        pw = 1.0
        add_round_rect(slide, X + W - pw - 0.26, y + 0.24, pw, 0.3, 'E7F5EC', 'BBE3C8', radius_in=0.15)
        add_text(slide, X + W - pw - 0.26, y + 0.24, pw, 0.3, d, size=10, bold=True, color='15924F', align='c', anchor='m')
    # accent tile
    ty = 2.54 + 3 * 0.90 + 0.06   # 5.30
    th = 6.46 - ty
    add_round_rect(slide, X, ty, W, th, '0F2A1E', '1F5A3B', radius_in=0.1)
    add_pic(slide, str(ICONDIR / 'clock_green.png'), X + 0.32, ty + 0.24, 0.34, 0.34)
    T = {'ru': dict(big='158 ч/мес', lab='возвращённого рабочего времени · Год 1',
                    pill='28 % рутины — у системы'),
         'en': dict(big='158 h/mo', lab='of working time returned · Year 1',
                    pill='28% of routine automated')}[L]
    add_text(slide, X + 0.82, ty + 0.10, 3.4, 0.72, T['big'], size=37, bold=True, color='3FB572', anchor='m')
    add_text(slide, X + 0.34, ty + 0.80, W - 0.6, 0.3, T['lab'], size=11.5, color='CFE0F1', anchor='m')
    pw = 2.05
    add_round_rect(slide, X + W - pw - 0.26, ty + 0.28, pw, 0.32, '10301F', '1D4A32', radius_in=0.16)
    add_text(slide, X + W - pw - 0.26, ty + 0.28, pw, 0.32, T['pill'], size=10, bold=True, color='3FB572', align='c', anchor='m')

# ============================================================
# move a slide within the deck
# ============================================================
def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_idx]
    sldIdLst.remove(el)
    sldIdLst.insert(to_idx, el)

def replace_text_everywhere(slide, old, new):
    for el in all_graphics(slide):
        for t in el.iter(q('a:t')):
            if t.text and old in t.text:
                t.text = t.text.replace(old, new)

# ============================================================
# DUPLICATE SLIDE (within same presentation)
# ============================================================
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# ---------------------------------------------------------------------------
# Repository layout: this file lives at <repo>/deck/build/transform.py
# ---------------------------------------------------------------------------
from pathlib import Path
ROOT   = Path(__file__).resolve().parents[2]
ASSETS = ROOT / 'assets'
ICONDIR= ASSETS / 'icons'
SHOTS  = ASSETS / 'shots'
BASE   = ROOT / 'deck' / 'base'
DECK   = ROOT / 'deck'


GRAPHIC_TAGS = None  # initialised below

def duplicate_slide(prs, index, insert_at):
    """Copy slide content into a new slide, keeping dest's live spTree object
    so python-pptx shape builders keep working on the duplicate."""
    tags = (q('p:sp'), q('p:pic'), q('p:graphicFrame'), q('p:grpSp'), q('p:cxnSp'))
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    # background
    src_cSld = src.part._element.find(q('p:cSld'))
    dest_cSld = dest.part._element.find(q('p:cSld'))
    src_bg = src_cSld.find(q('p:bg'))
    if src_bg is not None:
        old_bg = dest_cSld.find(q('p:bg'))
        if old_bg is not None: dest_cSld.remove(old_bg)
        dest_cSld.insert(0, copy.deepcopy(src_bg))
    # shapes: clear dest placeholders, copy src graphics into the LIVE dest tree
    dest_tree = dest.shapes._spTree
    src_tree = src_cSld.find(q('p:spTree'))
    for el in list(dest_tree):
        if el.tag in tags: dest_tree.remove(el)
    for el in src_tree:
        if el.tag in tags: dest_tree.append(copy.deepcopy(el))
    dest_el = dest.part._element
    # rebuild rels with rId mapping
    mapping = {}
    for rId, rel in list(src.part.rels.items()):
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE): continue
        if rel.is_external:
            newid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            newid = dest.part.relate_to(rel.target_part, rel.reltype)
        mapping[rId] = newid
    for el in dest_el.iter():
        for attr, val in list(el.attrib.items()):
            if val in mapping and ('}' in attr):
                el.set(attr, mapping[val])
    # move to insert_at
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    newId = ids[-1]
    sldIdLst.remove(newId)
    sldIdLst.insert(insert_at, newId)

def retitle(slide, kick_old_startswith, kick_new, title_old_startswith, title_new):
    # read spTree from the live part XML (slide.shapes may be a stale lazyproperty
    # for slides whose <p:sld> children were replaced wholesale)
    tree = slide.part._element.find(q('p:cSld')).find(q('p:spTree'))
    els = [el for el in tree if el.tag in (q('p:sp'), q('p:pic'))]
    for el in els:
        txt = sp_texts(el)
        if txt.startswith(kick_old_startswith):
            set_single_text(el, kick_new)
        elif txt.startswith(title_old_startswith) and 'BPMN' not in txt:
            set_single_text(el, title_new)

# ============================================================
# KZPromotion slide — the current tool and its limits
# ============================================================
def build_kzpromotion(slide, L):
    strip_slide(slide)
    S = {'ru': dict(
            kick='Текущий инструмент', h2='Текущая система KZPromotion — и её ограничения',
            sub='Сегодня уведомления партнёрам по акциям N+1 рассылает KZPromotion. Система закрывает лишь малую часть процесса и требует постоянного ручного участия.',
            cap='KZPromotion — рассылка уведомлений партнёрам об акциях N+1',
            lims=[('Только N+1', 'Остальные типы промо-активностей не поддерживаются.'),
                  ('Ручной выбор акций', 'Специалист Sales Systems отбирает промо-активности для рассылки вручную.'),
                  ('Автоматически — только квота', 'Уведомления о смене периода или продукта приходится отправлять партнёрам вручную.')],
            tt='Сбои в эксплуатации',
            tr=['Microsoft изменил внутренние политики — письма перестали отправляться, уведомления пришлось рассылать вручную.',
                'Мелкие сбои возникают регулярно.']),
         'en': dict(
            kick='Current tooling', h2='The current KZPromotion system — and its limits',
            sub='Today partner notifications for N+1 promos are sent by KZPromotion. It covers only a small part of the process and still requires constant manual work.',
            cap='KZPromotion — partner notification mailing for N+1 promos',
            lims=[('N+1 only', 'All other promo activity types are not supported.'),
                  ('Manual selection', 'The Sales Systems communication specialist picks promo activities for mailing by hand.'),
                  ('Quota changes only — automatic', 'Notifications about period or product changes have to be sent to partners manually.')],
            tt='Operational issues',
            tr=['Microsoft changed its internal policies — emails stopped going out, and notifications had to be sent manually.',
                'Minor failures occur frequently.'])}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_text(slide, MX, 1.56, 11.9, 0.6, S['sub'], size=12.5, color='56616F', spacing=1.15)
    top, bot = 2.32, 6.46
    # left: screenshot card + troubles card
    lw = 7.2
    img_w = lw - 0.48; img_h = img_w / (1280 / 378)   # ≈ 1.98
    sc_h = img_h + 0.66
    add_round_rect(slide, MX, top, lw, sc_h, 'FFFFFF', 'E2E7EE', shadow=True)
    pic = add_pic(slide, str(ASSETS / 'kzpromotion.png'), MX + 0.24, top + 0.2, img_w, img_h)
    pic.line.color.rgb = C('E2E7EE'); pic.line.width = Pt(0.75)
    add_text(slide, MX + 0.24, top + 0.24 + img_h, img_w, 0.34, S['cap'],
             size=10, bold=True, color='56616F', align='c', anchor='m')
    # troubles card (red tint)
    ty = top + sc_h + 0.18; th = bot - ty
    add_round_rect(slide, MX, ty, lw, th, 'FBE9E7', 'F0C5C0', radius_in=0.1)
    add_text(slide, MX + 0.3, ty + 0.12, lw - 0.6, 0.28, S['tt'].upper(), size=10, bold=True,
             color='C33A30', char_spacing=1, anchor='m')
    yy = ty + 0.44
    for t in S['tr']:
        add_rect(slide, MX + 0.32, yy + 0.09, 0.07, 0.07, 'C33A30')
        add_text(slide, MX + 0.52, yy, lw - 0.84, 0.42, t, size=10.5, color='0C1726', spacing=1.1)
        yy += 0.5
    # right: limitation cards
    rx = MX + lw + 0.18; rw = CW - lw - 0.18
    ch = (bot - top - 2 * 0.15) / 3
    for i, (t, b) in enumerate(S['lims']):
        y = top + i * (ch + 0.15)
        add_round_rect(slide, rx, y, rw, ch, 'FFFFFF', 'E2E7EE', shadow=True)
        add_rect(slide, rx, y + 0.02, 0.05, ch - 0.04, 'BD7D12')
        add_text(slide, rx + 0.26, y + 0.12, rw - 0.5, 0.3, t, size=12.5, bold=True, color='0C1726', anchor='m')
        add_text(slide, rx + 0.26, y + 0.46, rw - 0.5, ch - 0.56, b, size=10.5, color='56616F', spacing=1.12)
    # footer right text
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g and g[1] > 6.9 * EMU and 'BPMN' in sp_texts(el):
            set_single_text(el, 'Efes Nexus')

# ============================================================
# 6) video slide (replaces the screenshots on "Промо-дашборд")
# ============================================================
def build_video(slide, L):
    strip_slide(slide)
    S = {'ru': dict(kick='Система', h2='Промо-дашборд',
                    quote='«Единое окно для всего жизненного цикла промо-активности — от заявки партнёра до расчёта компенсации.»',
                    cap='Живая запись работы платформы'),
         'en': dict(kick='The system', h2='The promo dashboard',
                    quote='“A single window for the entire promo-activity lifecycle — from the partner’s request to the compensation calculation.”',
                    cap='Live recording of the platform')}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_pic(slide, str(ICONDIR / 'quote.png'), MX + 0.02, 1.70, 0.24, 0.24)
    add_text(slide, MX + 0.4, 1.62, CW - 0.5, 0.6, S['quote'], size=14.5, italic=True,
             color='13203A', spacing=1.2)
    # centered video card
    ASPECT = 3456 / 1840
    vh = 3.82; vw = vh * ASPECT               # ~7.18
    cx = (13.333 - vw) / 2
    card_y = 2.36
    add_round_rect(slide, cx - 0.12, card_y, vw + 0.24, vh + 0.62, 'FFFFFF', 'E2E7EE', shadow=True)
    mov = slide.shapes.add_movie(str(ASSETS / 'demo-video.mp4'), Inches(cx), Inches(card_y + 0.12),
                                 Inches(vw), Inches(vh),
                                 poster_frame_image=str(ASSETS / 'demo-poster.png'),
                                 mime_type='video/mp4')
    add_text(slide, cx, card_y + vh + 0.20, vw, 0.32, S['cap'],
             size=11, bold=True, color='56616F', align='c', anchor='m')

# ============================================================
# 8) logo replacement (badge -> new wordmark, keep width, center v)
# ============================================================
BADGE_BYTES = open(str(ASSETS / 'efes-badge.png'), 'rb').read()
LOGO_PATH = str(ASSETS / 'logo-efes.png')

def replace_logo(prs):
    from PIL import Image
    logo_bytes = open(LOGO_PATH, 'rb').read()
    w0, h0 = Image.open(LOGO_PATH).size
    aspect = w0 / h0
    parts = set()
    for slide in prs.slides:
        tree = slide.part._element.find(q('p:cSld')).find(q('p:spTree'))
        for pic in tree.iter(q('p:pic')):
            blip = pic.find('.//' + q('a:blip'))
            if blip is None: continue
            rId = blip.get(q('r:embed'))
            if rId is None: continue
            try:
                part = slide.part.related_part(rId)
            except KeyError:
                continue
            if part.blob == BADGE_BYTES:
                parts.add(part)
                g = sp_geom(pic)
                if g:
                    x, y, w, h = g
                    new_h = int(w / aspect)
                    off = pic.find('.//' + q('a:off')); ext = pic.find('.//' + q('a:ext'))
                    off.set('y', str(y + (h - new_h) // 2))
                    ext.set('cy', str(new_h))
    for part in parts:
        part._blob = logo_bytes

# ============================================================
# AUTHOR / TEAM SLIDE
# ============================================================
def plain_footer(slide):
    """Slides duplicated off a BPMN slide inherit its 'Efes Nexus · BPMN' footer."""
    for el in all_graphics(slide):
        g = sp_geom(el)
        if g and g[1] > 6.9 * EMU and 'BPMN' in sp_texts(el):
            set_single_text(el, 'Efes Nexus')

def author_photo():
    """assets/author.{jpg,jpeg,png} if the author dropped a photo in, else None."""
    for name in ('author.jpg', 'author.jpeg', 'author.png', 'author.webp'):
        p = ASSETS / name
        if p.exists():
            return str(p)
    return None

def build_author(slide, L):
    strip_slide(slide); plain_footer(slide)
    S = {'ru': dict(
            kick='Команда', h2='Команда проекта',
            sub='Проект выполняет команда из 5 инженеров — под руководством основателя проекта, '
                'отвечающего за архитектуру, ключевые технические решения и качество результата.',
            name='Асет Сексенали', initials='АС',
            role='Основатель проекта · Тимлид',
            d1='Senior Full-Stack инженер · 8+ лет высоконагруженных production-систем.',
            d2='Архитектура · решения · приёмка качества',
            cards=[('Halyk Bank',
                    'онлайн-банкинг: 100 000+ бизнес-клиентов, 20+ поставленных фич'),
                   ('Archimedes · CTO',
                    'общенациональная медицинская платформа — до 5 000 RPS, 5 регионов'),
                   ('Zencoder · США',
                    'AI-инструменты для разработчиков: плагины VS Code и JetBrains'),
                   ('Lumica · CTO',
                    'LMS-платформа: 15+ микросервисов, до 3 000 RPS, команда из 5 инженеров')],
            bandt='Команда: тимлид + 4 разработчика',
            bandb='backend · frontend · DevOps · QA — единый центр ответственности за результат'),
         'en': dict(
            kick='The team', h2='The project team',
            sub='The project is delivered by a team of 5 engineers — led by the founder of the project, '
                'who owns the architecture, the key technical decisions and the quality of the result.',
            name='Aset Seksenali', initials='AS',
            role='Project founder · Team lead',
            d1='Senior full-stack engineer · 8+ years of high-load production systems.',
            d2='Architecture · decisions · quality acceptance',
            cards=[('Halyk Bank',
                    'online banking: 100,000+ business clients, 20+ features shipped'),
                   ('Archimedes · CTO',
                    'nationwide medical platform — up to 5,000 RPS, 5 regions'),
                   ('Zencoder · USA',
                    'AI tools for developers: VS Code and JetBrains plugins'),
                   ('Lumica · CTO',
                    'LMS platform: 15+ microservices, up to 3,000 RPS, a team of 5 engineers')],
            bandt='The team: team lead + 4 developers',
            bandb='backend · frontend · DevOps · QA — a single point of accountability for the result')}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_text(slide, MX, 1.52, 11.9, 0.7, S['sub'], size=12.5, color='56616F', spacing=1.15)
    top, bot = 2.36, 6.46
    # ---- profile card (left) ----
    pw = 4.42; ph = bot - top
    add_round_rect(slide, MX, top, pw, ph, 'FFFFFF', 'E2E7EE', shadow=True)
    add_rect(slide, MX + 0.02, top, pw - 0.04, 0.05, '004C8D')
    av = 1.62; ax = MX + (pw - av) / 2; ay = top + 0.34
    photo = author_photo()
    if photo:
        add_pic(slide, photo, ax, ay, av, av)
    else:
        add_round_rect(slide, ax, ay, av, av, 'E3EEF8', '9CC6E8', radius_in=av / 2)
        add_text(slide, ax, ay, av, av, S['initials'], size=46, bold=True,
                 color='004C8D', align='c', anchor='m')
    add_text(slide, MX + 0.30, ay + av + 0.24, pw - 0.6, 0.42, S['name'], size=21, bold=True,
             color='0C1726', align='c', anchor='m')
    add_text(slide, MX + 0.30, ay + av + 0.70, pw - 0.6, 0.32, S['role'], size=12.5, bold=True,
             color='1166A8', align='c', anchor='m')
    add_rect(slide, MX + 0.85, ay + av + 1.18, pw - 1.7, 0.014, 'E2E7EE')
    add_text(slide, MX + 0.42, ay + av + 1.30, pw - 0.84, 0.56, S['d1'], size=11,
             color='56616F', align='c', spacing=1.25)
    add_text(slide, MX + 0.42, ay + av + 1.86, pw - 0.84, 0.30, S['d2'], size=11,
             color='56616F', align='c', anchor='m')
    # ---- track-record cards (right, 2x2) + team band ----
    rx = MX + pw + 0.24; rw = CW - pw - 0.24
    bh = 0.92                      # team band height
    gh = 0.22                      # gap
    ch = (ph - bh - gh - gh) / 2   # card row height
    cwr = (rw - gh) / 2
    for i, (lab, body) in enumerate(S['cards']):
        x = rx + (i % 2) * (cwr + gh)
        y = top + (i // 2) * (ch + gh)
        add_round_rect(slide, x, y, cwr, ch, 'FFFFFF', 'E2E7EE', shadow=True)
        add_text(slide, x + 0.30, y + 0.20, cwr - 0.58, 0.30, lab, size=13, bold=True,
                 color='0C1726', anchor='m')
        add_text(slide, x + 0.30, y + 0.56, cwr - 0.58, ch - 0.70, body, size=11,
                 color='56616F', spacing=1.25)
    by = top + 2 * (ch + gh)
    add_round_rect(slide, rx, by, rw, bh, 'E9F1FA', '9CC6E8')
    add_text(slide, rx + 0.30, by + 0.16, rw - 0.6, 0.30, S['bandt'], size=13, bold=True,
             color='004C8D', anchor='m')
    add_text(slide, rx + 0.30, by + 0.50, rw - 0.6, 0.30, S['bandb'], size=11,
             color='56616F', anchor='m')

# ============================================================
# COMMERCIAL PROPOSAL / PRICING SLIDE
# ============================================================
def build_pricing(slide, L):
    strip_slide(slide); plain_footer(slide)
    S = {'ru': dict(
            kick='Стоимость', h2='Стоимость проекта',
            sub='Фиксированная стоимость в два этапа: Этап 1 — разработка и внедрение системы; '
                'Этап 2 — опциональный пакет развития. Размещение — на инфраструктуре Efes, без затрат на хостинг.',
            l1='ЭТАП 1 · РАЗРАБОТКА И ВНЕДРЕНИЕ', p1='5 000 000 ₸',
            c1='фиксированная стоимость · без НДС',
            i1=[('8 модулей системы и 3 сквозных BPMN-процесса', None),
                ('Интеграции: Panorama, 1С, корпоративная почта', None),
                ('AI-верификация запросов партнёров', None),
                ('Развёртывание на сервере Efes и обучение команды', None),
                ('MVP — уже через 6 недель', None)],
            l2='ЭТАП 2 · РАЗВИТИЕ ПЛАТФОРМЫ (ОПЦИОНАЛЬНО)', p2='+ 5 000 000 ₸',
            c2='отдельное соглашение · модули в любом порядке',
            i2=[('AI Search & Knowledge Base', '1 600 000 ₸'),
                ('Проверка накладных 1С', '1 500 000 ₸'),
                ('Масштабирование расчётов (Kega Boom)', '1 100 000 ₸'),
                ('Исторический календарь промо', '800 000 ₸')]),
         'en': dict(
            kick='Pricing', h2='Project cost',
            sub='A fixed price in two phases: Phase 1 — development and rollout of the system; '
                'Phase 2 — an optional development package. Hosted on Efes infrastructure, with no hosting costs.',
            l1='PHASE 1 · DEVELOPMENT AND ROLLOUT', p1='5,000,000 KZT',
            c1='fixed price · VAT not included',
            i1=[('8 system modules and 3 end-to-end BPMN processes', None),
                ('Integrations: Panorama, 1C, corporate email', None),
                ('AI verification of partner requests', None),
                ('Deployment on the Efes server and team training', None),
                ('MVP — in just 6 weeks', None)],
            l2='PHASE 2 · PLATFORM DEVELOPMENT (OPTIONAL)', p2='+ 5,000,000 KZT',
            c2='separate agreement · modules in any order',
            i2=[('AI Search & Knowledge Base', '1,600,000 KZT'),
                ('1C invoice validation', '1,500,000 KZT'),
                ('Scaling the calculations (Kega Boom)', '1,100,000 KZT'),
                ('Historical promo calendar', '800,000 KZT')])}[L]
    kicker(slide, MX, 0.56, S['kick'], True)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='EEF2F7', anchor='m')
    add_text(slide, MX, 1.52, 11.9, 0.7, S['sub'], size=12.5, color='9AA9BF', spacing=1.15)
    top = 2.36; ch = 4.10
    cw = (CW - 0.30) / 2

    def card(x, lab, price, cap, items, primary):
        add_round_rect(slide, x, top, cw, ch, '16223B', '2E72AE' if primary else '28364E')
        add_text(slide, x + 0.46, top + 0.28, cw - 0.9, 0.28, lab, size=9, bold=True,
                 color='5CA3DB' if primary else '8FA0B8', char_spacing=1.2, anchor='m')
        add_text(slide, x + 0.46, top + 0.66, cw - 0.9, 0.86, price, size=40, bold=True,
                 color='FFFFFF', anchor='m')
        add_text(slide, x + 0.48, top + 1.58, cw - 0.9, 0.28, cap, size=11.5, italic=True,
                 color='5CA3DB' if primary else '8FA0B8', anchor='m')
        add_rect(slide, x + 0.46, top + 2.02, cw - 0.92, 0.014, '28364E')
        for j, (it, money) in enumerate(items):
            y = top + 2.18 + j * 0.38
            add_round_rect(slide, x + 0.48, y + 0.13, 0.075, 0.075,
                           '5CA3DB' if primary else '8FA0B8', radius_in=0.0375)
            runs = [(it, {})] if money is None else [(it + ' — ', {}),
                                                     (money, {'bold': True, 'color': 'FFFFFF'})]
            add_text(slide, x + 0.74, y, cw - 1.2, 0.32, runs, size=11.5,
                     color='C8D3E2', anchor='m')

    card(MX, S['l1'], S['p1'], S['c1'], S['i1'], True)
    card(MX + cw + 0.30, S['l2'], S['p2'], S['c2'], S['i2'], False)

# ============================================================
# DELETE SLIDE
# ============================================================
def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldId = slides[index]
    rId = sldId.get(q('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# ============================================================
# 1C integration — matching the universal report against the invoices
# ============================================================
def build_1c(slide, L):
    strip_slide(slide); plain_footer(slide)
    S = {'ru': dict(
            kick='Интеграция · 1С', h2='Интеграция с 1С — сверка отчёта и накладных',
            sub='Партнёр присылает выгрузку из 1С и накладные — система сверяет их между собой и '
                'подсвечивает расхождения. Проверку становится возможно проводить сразу по всем '
                'торговым точкам: вручную такой объём данных не обработать.',
            cap1='Универсальный отчёт из 1С',
            cap2='Накладная на отпуск запасов на сторону',
            chip='СВЕРКА',
            cards=[('Загрузка отчёта из 1С',
                    'Универсальный отчёт по реализации загружается одним файлом.'),
                   ('Сверка с накладными',
                    'Суммы, скидки и объёмы сопоставляются по каждому контрагенту.'),
                   ('Контроль расхождений',
                    'Несовпадения по скидке, объёму или периоду подсвечиваются автоматически.'),
                   ('Все точки сразу',
                    'Можно запросить накладные со всех точек партнёра — вручную это невозможно.')]),
         'en': dict(
            kick='Integration · 1C', h2='1C integration — matching the report against the invoices',
            sub='The partner sends the 1C export and the invoices — the system matches them against '
                'each other and highlights the discrepancies. Checking becomes possible across every '
                'outlet at once: that volume of data cannot be processed by hand.',
            cap1='Universal report from 1C',
            cap2='Goods release note (invoice)',
            chip='MATCHING',
            cards=[('1C report upload',
                    'The universal sales report is uploaded as a single file.'),
                   ('Matching against invoices',
                    'Amounts, discounts and volumes are matched per counterparty.'),
                   ('Discrepancy control',
                    'Mismatches in discount, volume or period are highlighted automatically.'),
                   ('Every outlet at once',
                    'Invoices can be requested from all partner outlets — impossible by hand.')])}[L]
    kicker(slide, MX, 0.56, S['kick'], False)
    add_text(slide, MX, 0.86, CW, 0.66, S['h2'], size=29, bold=True, color='0C1726', anchor='m')
    add_text(slide, MX, 1.52, 11.9, 0.72, S['sub'], size=12.5, color='56616F', spacing=1.15)

    top = 2.42
    ih = 1.86                                   # shared display height of both screenshots
    w1 = ih * (871 / 528); w2 = ih * (1280 / 387)
    pad = 0.22
    c1w = w1 + pad * 2; c2w = w2 + pad * 2
    gap = CW - c1w - c2w
    ch = ih + pad * 2 + 0.32                    # + caption strip
    x1 = MX; x2 = MX + c1w + gap

    def shot(x, cw_, path, cap):
        add_round_rect(slide, x, top, cw_, ch, 'FFFFFF', 'E2E7EE', shadow=True)
        pic = add_pic(slide, path, x + pad, top + pad, cw_ - pad * 2, ih)
        pic.line.color.rgb = C('E2E7EE'); pic.line.width = Pt(0.75)
        add_text(slide, x + pad, top + pad + ih, cw_ - pad * 2, 0.32, cap,
                 size=10, bold=True, color='56616F', align='c', anchor='m')

    shot(x1, c1w, str(SHOTS / '1c_report.png'), S['cap1'])
    shot(x2, c2w, str(SHOTS / '1c_invoice.png'), S['cap2'])

    # centre chip between the two screenshots
    cd = 0.72; cx = x1 + c1w + (gap - cd) / 2; cy = top + (ch - cd) / 2 - 0.16
    add_round_rect(slide, cx, cy, cd, cd, '004C8D', None, radius_in=cd / 2)
    add_pic(slide, str(ICONDIR / 'invoice_white.png'), cx + cd * 0.26, cy + cd * 0.26,
            cd * 0.48, cd * 0.48)
    add_text(slide, cx - 0.34, cy + cd + 0.08, cd + 0.68, 0.26, S['chip'], size=8.5, bold=True,
             color='004C8D', char_spacing=1.2, align='c', anchor='m')

    # functionality row
    by = top + ch + 0.20; bh = 6.46 - by
    g = 0.18; cw_ = (CW - g * 3) / 4
    for i, (t, b) in enumerate(S['cards']):
        x = MX + i * (cw_ + g)
        add_round_rect(slide, x, by, cw_, bh, 'FFFFFF', 'E2E7EE', shadow=True)
        add_rect(slide, x, by + 0.02, 0.05, bh - 0.04, '004C8D')
        add_text(slide, x + 0.24, by + 0.12, cw_ - 0.46, 0.28, t, size=12, bold=True,
                 color='0C1726', anchor='m')
        add_text(slide, x + 0.24, by + 0.44, cw_ - 0.46, bh - 0.54, b, size=10,
                 color='56616F', spacing=1.14)

# ============================================================
def run(src, out, L, shots):
    prs = Presentation(src)
    # keep a pristine copy of BPMN slide as new slide 3 ("with the system")
    duplicate_slide(prs, 1, 2)
    transform_slide2(prs.slides[1], L)
    T = {'ru': dict(k2='BPMN · КАК СЕЙЧАС — БЕЗ СИСТЕМЫ',
                    t2='Жизненный цикл промо-активности — сегодня',
                    k3='BPMN · ЧТО АВТОМАТИЗИРУЕТ СИСТЕМА',
                    t3='Жизненный цикл промо-активности — с Nexus',
                    told='Жизненный цикл'),
         'en': dict(k2='BPMN · TODAY — WITHOUT THE SYSTEM',
                    t2='Promo activity lifecycle — today',
                    k3='BPMN · WHAT THE SYSTEM AUTOMATES',
                    t3='Promo activity lifecycle — with Nexus',
                    told='Promo activity lifecycle')}[L]
    retitle(prs.slides[1], 'BPMN ·', T['k2'], T['told'], T['t2'])
    retitle(prs.slides[2], 'BPMN ·', T['k3'], T['told'], T['t3'])
    remove_red_node(prs.slides[1], L)          # (1) drop the not-performed step from "today"
    # today the partner letter is written and sent by hand, so it is not an auto-notification
    MAIL = {'ru': ('Авто-уведомление партнёров письмом', 'Ручное уведомление партнёров письмом'),
            'en': ('Auto-notifying partners by email', 'Manual partner notification by email')}[L]
    replace_text_everywhere(prs.slides[1], *MAIL)
    mark_new_step(prs.slides[2], L)            # (3+10) red NEW step + accent callout
    build_video(prs.slides[4], L)              # (5+6) "Промо-дашборд" with embedded video
    stats_v3(prs.slides[5], L)                 # no charts; saved time is the accent
    build_future(prs.slides[6], L)             # (7+9)
    # (4) title slide rename -> Promo Management Platform
    replace_text_everywhere(prs.slides[0], 'Менеджер промо-активностей', 'Платформа управления промо-активностями')
    replace_text_everywhere(prs.slides[0], 'Promo Activity Manager', 'Promo Management Platform')
    # (2) stats slide -> position 3 (index 2)
    move_slide(prs, 5, 2)
    # KZPromotion (current tool) slide right after "today" -> index 2, stats shifts to 3
    duplicate_slide(prs, 1, 2)
    build_kzpromotion(prs.slides[2], L)
    # closing slides: who builds it, and what it costs. Appended before the delete
    # below, otherwise they would reuse the dropped slide's part name.
    n = len(prs.slides._sldIdLst)
    duplicate_slide(prs, 1, n)
    build_author(prs.slides[n], L)
    # pricing is duplicated off the dark 'Этап 2' slide (index 7) so it inherits the dark
    # background — the two commercial slides then read as one block.
    duplicate_slide(prs, 7, n + 1)
    build_pricing(prs.slides[n + 1], L)
    # closing deep-dive: how the 1C integration actually checks the data
    duplicate_slide(prs, 1, n + 2)
    build_1c(prs.slides[n + 2], L)
    # drop the old 'scaling' slide last (its part number would otherwise be reused)
    delete_slide(prs, 8)
    # (8) new logo everywhere
    replace_logo(prs)
    prs.save(out)
    print('saved', out, '| slides:', len(prs.slides._sldIdLst))

if __name__ == '__main__':
    run(str(BASE / 'Efes_Nexus_RU_base.pptx'), str(DECK / 'Efes_Nexus_RU.pptx'), 'ru',
        [str(SHOTS / 'ru_list.png'), str(SHOTS / 'ru_detail.png')])
    run(str(BASE / 'Efes_Nexus_EN_base.pptx'), str(DECK / 'Efes_Nexus_EN.pptx'), 'en',
        [str(SHOTS / 'en_list.png'), str(SHOTS / 'en_detail.png')])
